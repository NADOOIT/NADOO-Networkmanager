from pptx import Presentation
from sqlalchemy.orm import Session
from data.models import PresentationSlide
from pptx.dml.color import RGBColor
from pptx.util import Pt


class SlideTemplate:
    def __init__(self, db_session: Session, pptx_path: str):
        self.db_session = db_session
        self.pptx_path = pptx_path

    def fetch_slide_data(self, slide_id: int):
        return self.db_session.query(PresentationSlide).filter(PresentationSlide.ID == slide_id).first()

    def create_from_template(self, slide_id: int, output_path: str):
        print("Creating presentation...")
        presentation_data = self.fetch_slide_data(slide_id)
        if not presentation_data:
            raise ValueError(f"No data found for slide ID {slide_id}")

        prs = Presentation(self.pptx_path)

        # Assuming the slide to update is the first slide in the presentation
        slide = prs.slides[0]

        # Iterate over the shapes and update text or image
        for shape in slide.shapes:
            if shape.name == "PhotoShape":  # Picture name from the .pptx doc: Home -> Arrange -> Selection Pane
                # Replace the picture
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                # Remove the old picture
                slide.shapes._spTree.remove(shape._element)
                # Add the new picture
                slide.shapes.add_picture("resources/images/Picture1.png", x, y, cx, cy)
                print("Picture shape updated successfully")

            elif shape.has_text_frame:
                text_frame = shape.text_frame
                text = shape.text

                if shape.text == "Kurzpräsentation":
                    text_frame.clear()  # not necessary for newly-created shape
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = presentation_data.Folientitel
                    font = run.font
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(27)
                    font.bold = False
                    font.italic = None  # cause value to be inherited from theme
                elif shape.text == "Firmenname":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = presentation_data.Firmenname
                    font = run.font
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(22)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Unternehmensbranche":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = presentation_data.Unternehmensbranche
                    font = run.font
                    font.color.rgb = RGBColor.from_string("FF0000")
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(20)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Kontaktdaten":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = presentation_data.Kontaktdaten
                    font = run.font
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(13)
                    font.bold = False
                    font.italic = None
                elif shape.text == "00 Sek":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = presentation_data.Vortragszeit
                    font = run.font
                    font.color.rgb = RGBColor.from_string("FFFFFF")
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(13)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Vorname":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = presentation_data.Vorname
                    font = run.font
                    font.color.rgb = RGBColor.from_string("FFFFFF")
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(24)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Nachname":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = presentation_data.Nachname
                    font = run.font
                    font.color.rgb = RGBColor.from_string("FFFFFF")
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(24)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Vor- und Nachname":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = presentation_data.Naechster_vortrag
                    font = run.font
                    font.name = 'Arial'
                    font.size = Pt(24)
                    font.bold = True
                    font.italic = None

            prs.save(output_path)
            print("Text frame updated successfully")
