from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import os


class FolienVorlage:
    def __init__(self, folienvorlage, user_info: dict = None):
        self.user_info = user_info
        self.folienvorlage = folienvorlage
        self.folientitel = folienvorlage['folientitel']
        self.pptx_src_path = folienvorlage['folien_path']
        self.pptx_dest_path = folienvorlage['folien_path']

    def create_from_template(self):
        print("Creating presentation...")

        prs = Presentation(self.pptx_src_path)

        # Assuming the slide to update is the first slide in the presentation
        slide = prs.slides[0]

        # Iterate over the shapes and update text or image
        for shape in slide.shapes:
            if shape.name == "PhotoShape":  # Picture name from the .pptx doc: Home -> Arrange -> Selection Pane
                # Replace the picture
                x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
                # Remove the old picture
                slide.shapes._spTree.remove(shape._element)
                # Add the new picture
                slide.shapes.add_picture("resources/images/user.png", x, y, cx, cy)
                print("Picture shape updated successfully")

            elif shape.has_text_frame:
                text_frame = shape.text_frame
                text = shape.text

                if shape.text == "Kurzpräsentation":
                    text_frame.clear()  # not necessary for newly-created shape
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = self.folientitel
                    font = run.font
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(27)
                    font.bold = False
                    font.italic = None  # cause value to be inherited from theme
                elif shape.text == "Firmenname":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = self.user_info['firmenname']
                    font = run.font
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(22)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Unternehmensbranche":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = self.user_info['unternehmensbranche']
                    font = run.font
                    font.color.rgb = RGBColor.from_string("FF0000")
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(20)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Kontaktdaten":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = \
                        f"Telefon: {self.user_info['telefonnummer']}\n{self.user_info['email']}\n{self.user_info['webseite']}"
                    font = run.font
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(13)
                    font.bold = False
                    font.italic = None
                elif shape.text == "00 Sek":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = self.folienvorlage['vortragszeit']
                    font = run.font
                    font.color.rgb = RGBColor.from_string("FFFFFF")
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(13)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Vorname":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = self.user_info['vorname']
                    font = run.font
                    font.color.rgb = RGBColor.from_string("FFFFFF")
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(24)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Nachname":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = self.user_info['nachname']
                    font = run.font
                    font.color.rgb = RGBColor.from_string("FFFFFF")
                    font.name = 'Arial Black (Headings)'
                    font.size = Pt(24)
                    font.bold = False
                    font.italic = None
                elif shape.text == "Vor- und Nachname":
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = "Max Mustermann"
                    font = run.font
                    font.name = 'Arial'
                    font.size = Pt(24)
                    font.bold = True
                    font.italic = None

        # Speichere die Änderungen in einer neuen Datei.
        directory, filename = os.path.split(self.pptx_dest_path)
        name, ext = os.path.splitext(filename)
        new_filename = f"{self.user_info['vorname']}_{self.folienvorlage['folientitel']}{ext}"
        new_directory = os.path.join(directory, "folien")
        if not os.path.exists(new_directory):
            os.makedirs(new_directory)
        destination_path = os.path.join(new_directory, new_filename)
        prs.save(destination_path)
        print("Text frame updated successfully")
