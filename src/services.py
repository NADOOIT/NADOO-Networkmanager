import os
import shutil
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from src.utils import ensure_directory_exists
from src.models.user import User
from src.validators import ValidationError, benutzerdaten_validieren
from src.data.storage import benutzer_speichern, get_benutzer_liste


def benutzerfoto_speichern(photo_path, vorname, nachname):
    if photo_path and photo_path != 'Kein Foto ausgewählt':
        photo_dir = 'resources/images/benutzer'
        os.makedirs(photo_dir, exist_ok=True)
        _, ext = os.path.splitext(photo_path)  # Extract file extension
        photo_filename = f'{vorname}_{nachname}_foto{ext}'
        photo_destination = os.path.join(photo_dir, photo_filename)
        photo_destination = photo_destination.replace('\\', '/')
        if os.path.abspath(photo_path) != os.path.abspath(photo_destination):
            shutil.copy(photo_path, photo_destination)
        return photo_destination
    return photo_path


def benutzerfoto_loeschen(user):
    user_photo_path = user.get('foto')
    if user_photo_path and os.path.exists(user_photo_path):
        os.remove(user_photo_path)


def benutzerdaten_validieren_und_speichern(benutzerdaten, data):
    # Validate the data
    benutzerdaten_validieren(benutzerdaten)

    if benutzerdaten['foto'] != 'Kein Foto ausgewählt':
        benutzerdaten['foto'] = benutzerfoto_speichern(benutzerdaten['foto'], benutzerdaten['vorname'],
                                                       benutzerdaten['nachname'])
    else:
        benutzerdaten['foto'] = 'resources/images/benutzer/user.png'

    # Data is valid, insert the data into the JSON file
    user = User(**benutzerdaten)

    # Find the highest current id
    if data['benutzer']:
        max_id = max(user['id'] for user in data['benutzer'])
    else:
        max_id = 0

    # Set the new user's id
    new_user_data = user.to_dict()
    new_user_data['id'] = max_id + 1

    data['benutzer'].append(new_user_data)
    benutzer_speichern(data)


def benutzerdaten_validieren_und_aktualisieren(benutzerdaten, data, benutzer_id):
    benutzerdaten_validieren(benutzerdaten)

    if benutzerdaten['foto'] != 'Kein Foto ausgewählt':
        benutzerdaten['foto'] = benutzerfoto_speichern(benutzerdaten['foto'], benutzerdaten['vorname'],
                                                       benutzerdaten['nachname'])
    elif benutzerdaten['foto'] == 'Kein Foto ausgewählt' or benutzerdaten['foto'] == '' or not benutzerdaten[
        'foto']:
        benutzerdaten['foto'] = 'resources/images/benutzer/user.png'


    # Find the user and update the information
    for benutzer in data['benutzer']:
        if benutzer['id'] == benutzer_id:
            benutzer['vorname'] = benutzerdaten['vorname']
            benutzer['nachname'] = benutzerdaten['nachname']
            benutzer['firmenname'] = benutzerdaten['firmenname']
            benutzer['unternehmensbranche'] = benutzerdaten['unternehmensbranche']
            benutzer['telefonnummer'] = benutzerdaten['telefonnummer']
            benutzer['email'] = benutzerdaten['email']
            benutzer['webseite'] = benutzerdaten['webseite']
            benutzer['chapter'] = benutzerdaten['chapter']
            benutzer['mitgliedsstatus'] = benutzerdaten['mitgliedsstatus']
            if benutzerdaten['foto'] != 'Kein Foto ausgewählt':
                benutzer['foto'] = benutzerdaten['foto']
            break

    benutzer_speichern(data)


def benutzerdaten_geaendert(user_form, selected_user) -> (bool, dict):
    benutzerdaten_aus_feldern = user_form.benutzerdaten_aus_feldern()
    benutzerdaten_aus_feldern['id'] = selected_user['id']

    benutzer_info = [user for user in get_benutzer_liste() if user['id'] == selected_user['id']][0]
    # True if the data is different
    return benutzer_info != benutzerdaten_aus_feldern, benutzerdaten_aus_feldern


def kurzpraesentation_folie_erzeugen(folienvorlage, user_info: dict = None):
    """
    Creates a new presentation from the template.
    :return: The path to the newly created presentation.
    """
    print("Folie wird erzeugt...")

    folienvorlage = folienvorlage
    folientitel = folienvorlage['folientitel']
    pptx_src_path = folienvorlage['folien_path']

    presentation = Presentation(pptx_src_path)

    # Assuming the slide to update is the first slide in the presentation
    folie = presentation.slides[0]

    # Iterate over the shapes and update text or image
    for shape in folie.shapes:
        if shape.name == "PhotoShape":  # Picture name from the .pptx doc: Home -> Arrange -> Selection Pane
            # Replace the picture
            x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
            # Remove the old picture
            folie.shapes._spTree.remove(shape._element)
            # Add the new picture
            folie.shapes.add_picture(user_info.get('foto', "resources/images/benutzer/user.png"), x, y, cx, cy)
            # print("Picture shape updated successfully")

        elif shape.has_text_frame:
            text_frame = shape.text_frame
            text = shape.text

            if shape.text == "Kurzpräsentation":
                text_frame.clear()  # not necessary for newly-created shape
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = folientitel
                font = run.font
                font.name = 'Arial Black (Headings)'
                font.size = Pt(27)
                font.bold = False
                font.italic = None  # cause value to be inherited from theme
                # print(f"Kurzpräsentation in {folientitel} geändert.")
            elif shape.text == "Firmenname":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = user_info['firmenname']
                font = run.font
                font.name = 'Arial Black (Headings)'
                font.size = Pt(22)
                font.bold = False
                font.italic = None
                # print(f"Firmenname in {user_info['firmenname']} geändert.")
            elif shape.text == "Unternehmensbranche":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = user_info['unternehmensbranche']
                font = run.font
                font.color.rgb = RGBColor.from_string("FF0000")
                font.name = 'Arial Black (Headings)'
                font.size = Pt(20)
                font.bold = False
                font.italic = None
                # print(f"Unternehmensbranche in {user_info['unternehmensbranche']} geändert.")
            elif shape.text == "Kontaktdaten":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = \
                    f"Telefon: {user_info['telefonnummer']}\n{user_info['email']}\n{user_info['webseite']}"
                font = run.font
                font.name = 'Arial Black (Headings)'
                font.size = Pt(13)
                font.bold = False
                font.italic = None
            elif shape.text == "Vorname":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = user_info['vorname']
                font = run.font
                font.color.rgb = RGBColor.from_string("FFFFFF")
                font.name = 'Arial Black (Headings)'
                font.size = Pt(24)
                font.bold = False
                font.italic = None
            elif shape.text == "Nachname":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = user_info['nachname']
                font = run.font
                font.color.rgb = RGBColor.from_string("FFFFFF")
                font.name = 'Arial Black (Headings)'
                font.size = Pt(24)
                font.bold = False
                font.italic = None
            elif shape.text == "Vor- und Nachname":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = "Max Mustermann"
                font = run.font
                font.name = 'Arial'
                font.size = Pt(24)
                font.bold = True
                font.italic = None

    dateizielpfad = kurzpraesentation_zielpfad_erstellen(user_info, folienvorlage)
    presentation.save(dateizielpfad)


def kurzpraesentation_zielpfad_erstellen(user_info, folienvorlage) -> str:
    """
    Erstellt einen Ordner für die Kurzpräsentationen.
    :return: Pfad zum Ordner
    """
    # Speichere die Änderungen in einer neuen Datei.
    pptx_src_path = folienvorlage['folien_path']
    ordner, dateiname = os.path.split(pptx_src_path)
    name, ext = os.path.splitext(dateiname)
    neuer_dateiname = f"{user_info['vorname']}_{folienvorlage['folientitel']}{ext}"
    neuer_ordner = os.path.join(ordner, "kurzpraesentation")
    ensure_directory_exists(neuer_ordner)
    dateizielpfad = os.path.join(neuer_ordner, neuer_dateiname)
    return dateizielpfad.replace('\\', '/')
