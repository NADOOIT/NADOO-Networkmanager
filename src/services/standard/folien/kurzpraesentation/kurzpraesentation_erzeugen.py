from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx import Presentation

from .kurzpraesentation_zielpfad_erstellen import kurzpraesentation_zielpfad_erstellen


def kurzpraesentation_erzeugen(folienvorlage, benutzer_folie, user_info: dict = None):
    """
    Creates a new presentation from the template.
    :return: The path to the newly created presentation.
    """

    folienvorlage = folienvorlage
    folientitel = folienvorlage['folientitel']
    pptx_src_path = folienvorlage['folien_path']

    print(benutzer_folie)

    presentation = Presentation(pptx_src_path)

    # Assuming the slide to update is the first slide in the presentation
    folie = presentation.slides[0]

    # Iterate over the shapes and update text or image
    for shape in folie.shapes:
        if shape.name == "PhotoShape":  # Picture name from the .pptx doc: Home -> Arrange -> Selection Pane
            # Replace the picture
            x, y, cx, cy = shape.left, shape.top, shape.width, shape.height
            # Remove the old picture
            folie.shapes._spTree.remove(shape._element)
            # Add the new picture
            folie.shapes.add_picture(user_info.get('foto', "resources/images/benutzer/user.png"), x, y, cx, cy)
            print(user_info.get('foto'))
            # print("Picture shape updated successfully")

        elif shape.has_text_frame:
            text_frame = shape.text_frame
            text = shape.text

            if shape.text == "Kurzpräsentation":
                text_frame.clear()  # not necessary for newly-created shape
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = folientitel
                font = run.font
                font.name = 'Arial Black (Headings)'
                font.size = Pt(27)
                font.bold = False
                font.italic = None  # cause value to be inherited from theme
                # print(f"Kurzpräsentation in {folientitel} geändert.")
            elif shape.text == "Firmenname":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = user_info['firmenname']
                font = run.font
                font.name = 'Arial Black (Headings)'
                font.size = Pt(22)
                font.bold = False
                font.italic = None
                # print(f"Firmenname in {user_info['firmenname']} geändert.")
            elif shape.text == "Unternehmensbranche":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = user_info['unternehmensbranche']
                font = run.font
                font.color.rgb = RGBColor.from_string("FF0000")
                font.name = 'Arial Black (Headings)'
                font.size = Pt(20)
                font.bold = False
                font.italic = None
                # print(f"Unternehmensbranche in {user_info['unternehmensbranche']} geändert.")
            elif shape.text == "Kontaktdaten":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = \
                    f"Telefon: {user_info['telefonnummer']}\n{user_info['email']}\n{user_info['webseite']}"
                font = run.font
                font.name = 'Arial Black (Headings)'
                font.size = Pt(13)
                font.bold = False
                font.italic = None
            elif shape.text == "Vorname":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = user_info['vorname']
                font = run.font
                font.color.rgb = RGBColor.from_string("FFFFFF")
                font.name = 'Arial Black (Headings)'
                font.size = Pt(24)
                font.bold = False
                font.italic = None
            elif shape.text == "Nachname":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = user_info['nachname']
                font = run.font
                font.color.rgb = RGBColor.from_string("FFFFFF")
                font.name = 'Arial Black (Headings)'
                font.size = Pt(24)
                font.bold = False
                font.italic = None
            elif shape.text == "Vor- und Nachname":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = benutzer_folie['naechster_vortrag'] or "Kein Vortrag"
                font = run.font
                font.name = 'Arial Black (Headings)'
                font.size = Pt(24)
                font.bold = False
                font.italic = None
            elif shape.text == "00 Sek":
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = benutzer_folie['vortragszeit']
                font = run.font
                font.color.rgb = RGBColor.from_string("FFFFFF")
                font.name = 'Arial'
                font.size = Pt(13)
                font.bold = True
                font.italic = None

            print(shape.text)

    dateizielpfad = kurzpraesentation_zielpfad_erstellen(user_info, folienvorlage)
    # print(f"Kurzpräsentation gespeichert in {dateizielpfad}")
    presentation.save(dateizielpfad)
